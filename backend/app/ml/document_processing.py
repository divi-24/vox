"""
Document processing service for PDFs, Word docs, PowerPoint slides.

Extracts text, decisions, and action items from various document types.
"""

import logging
from typing import List, Dict, Any
import asyncio
from pathlib import Path

import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Process various document types and extract structured intelligence"""

    async def process_document(self, file_path: str, doc_type: str) -> Dict[str, Any]:
        """
        Process document file.

        Args:
            file_path: Path to document
            doc_type: pdf, docx, pptx

        Returns:
            {
                "text": full extracted text,
                "pages": number of pages,
                "sections": [{"title": "...", "content": "..."}],
                "document_type": str
            }
        """

        processor = {
            "pdf": self._process_pdf,
            "docx": self._process_docx,
            "pptx": self._process_pptx,
        }.get(doc_type.lower())

        if not processor:
            raise ValueError(f"Unsupported document type: {doc_type}")

        return await asyncio.to_thread(processor, file_path)

    def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract text and structure from PDF"""

        try:
            sections = []
            all_text = []

            with open(file_path, "rb") as file:
                pdf = PyPDF2.PdfReader(file)
                total_pages = len(pdf.pages)

                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    all_text.append(text)

                    sections.append({
                        "page": page_num + 1,
                        "title": f"Page {page_num + 1}",
                        "content": text,
                    })

            full_text = "\n\n".join(all_text)

            logger.info(f"Processed PDF: {total_pages} pages")

            return {
                "text": full_text,
                "pages": total_pages,
                "sections": sections,
                "document_type": "pdf",
            }

        except Exception as e:
            logger.error(f"PDF processing failed: {e}")
            raise

    def _process_docx(self, file_path: str) -> Dict[str, Any]:
        """Extract text and structure from Word document"""

        try:
            sections = []
            all_text = []

            doc = DocxDocument(file_path)

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    all_text.append(paragraph.text)

                    # Detect section headers (usually bold or larger text)
                    if paragraph.style.name.startswith("Heading"):
                        sections.append({
                            "title": paragraph.text,
                            "content": "",
                            "level": int(paragraph.style.name[-1]) if paragraph.style.name[-1].isdigit() else 1,
                        })
                    elif sections:
                        sections[-1]["content"] += paragraph.text + "\n"

            full_text = "\n\n".join(all_text)

            logger.info(f"Processed DOCX: {len(doc.paragraphs)} paragraphs")

            return {
                "text": full_text,
                "pages": 1,
                "sections": sections,
                "document_type": "docx",
            }

        except Exception as e:
            logger.error(f"DOCX processing failed: {e}")
            raise

    def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """Extract text and structure from PowerPoint"""

        try:
            sections = []
            all_text = []

            prs = Presentation(file_path)

            for slide_num, slide in enumerate(prs.slides):
                slide_content = []

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                        all_text.append(shape.text)

                sections.append({
                    "slide": slide_num + 1,
                    "title": slide_content[0] if slide_content else f"Slide {slide_num + 1}",
                    "content": "\n".join(slide_content),
                })

            full_text = "\n\n".join(all_text)

            logger.info(f"Processed PPTX: {len(prs.slides)} slides")

            return {
                "text": full_text,
                "pages": len(prs.slides),
                "sections": sections,
                "document_type": "pptx",
            }

        except Exception as e:
            logger.error(f"PPTX processing failed: {e}")
            raise
